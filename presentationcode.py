from pptx import Presentation
from pptx.util import Inches, Pt
from gigachat import GigaChat
from gigachat.models import Chat, Messages, MessagesRole

def MyPermit(permit, contect):
    if type(permit) is str and type(contect) is str:
        payload = Chat(
            messages=[
                Messages(
                    role=MessagesRole.SYSTEM,
                    content=contect
                )
            ],
            temperature=0.7,
            max_tokens=100,
        )
        with GigaChat(credentials="ZTcwNWE2ZTYtMzNmNy00YTZmLThlMDMtMGE0NjdhYWJkMTY1OmUwZmQzNjY0LTNhYmUtNGRjZC1hMDJlLTM3M2JlOTk1NjRlNg==", verify_ssl_certs=False) as giga:
            payload.messages.append(Messages(role=MessagesRole.USER, content=permit))
            response = giga.chat(payload)
            payload.messages.append(response.choices[0].message)
            print(response.choices[0].message.content)
    return str(response.choices[0].message.content)

uslovia = ". Распиши это так, как будто ты делаешь доклад в очень строгой форме. Уложи это в менее 30 слов и одно предложение. Говори от третьего лица и не расписывай свои действия. Не принимай себя ни за мужчину, ни за женщину"


#Выбор языков
lang = input("Какой язык вы хотите использовать? (рус или англ) \n Wich language do you want to use? (rus or ing)").lower()
if lang == "rus" or lang == "рус":
    language = "Напиши это на русском"
    firstTemplateslide = "История"
    secondTemplateslide = "Основные понятия"
    thirdTemplateslide = "Предаставление знаний"
    fourthTemplateslide = "Задачи формирования баз знаний"
    fifthTemplateslide = "Заключение"
    qwsb = "Введите тему"
    plos = "Введите место обучения"
    doc = "Введите кафедру или класс обучения"
    aut = "Введите Ваше имя"
    mng = "Введите имя научного руководителя. Если его нет, не пишите ничего"
    plc = "Введите город, где защищаетесь"
    year = "Введите год защиты проекта"
    vot =  "У Вас есть свой собственный шаблон презентации? Напишите: \"да\" или \"нет\""
    nos = "Сколько Вам нужно слайдов? Введите число"
    nmos = "Введите название следующего слайда"
elif lang == "ing" or lang == "англ":
    language = "Write it in English"
    firstTemplateslide = "History"
    secondTemplateslide = "Basic concepts of "
    thirdTemplateslide = "Representation knowledge"
    fourthTemplateslide = "Objectives"
    fifthTemplateslide = "Conclusion"
    qwsb = "input subject"
    plos = "input your place of study"
    doc = "input your academic departament or academic class"
    aut = "input your name and surname"
    mng = "input your manager's name and surname"
    plc = "input your location (city or another settlement)"
    year = "input year of presentation"
    vot = "Do you have your own template for this presentation? Write: \"yes\" or \"no\""
    nos = "How many slides do you need? Enter a number"
    nmos = "Input name of next slide"
#TittleList
subject= input(qwsb)
placeofstudy= input(plos)
departmentorclass= input(doc)
autor= input(aut)
manager= input(mng)
place= input(plc)
year= input(year)

slides = []

#Создаем объект Presentation
presentation = Presentation()

#CustomTemplate
VariantOfTemplate = input(vot)
if VariantOfTemplate.lower() == "yes" or VariantOfTemplate.lower() == "да":
    numberofslides = int(input(nos))
    for i in range (numberofslides):
            slides.append(input(nmos))

    #TittleSlide
    slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(slide_layout)

    left = top = Inches(0)
    height = Inches(7)
    width = Inches(10)
    #pic = slide.shapes.add_picture(img_path1, left, top, height=height, width=width)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = subject
    subtitle.text = "" + placeofstudy + "\n" \
                    "" + departmentorclass + "\n" \
                    '' + autor + '\n' \
                    "" + manager + "\n" \
                    "" + place + ", " + year + ""
    number = 0
    for i in range (numberofslides):
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slides[number]

        content.text = (MyPermit("Напиши "+slides[number]+" объекта (вещи) "+ subject,uslovia +language)) #на это должен отвечать gigachad в виде: write [Notion of SUBJECT(НАЗВАНИЕ СЛАЙДА)]

        number+=1


elif VariantOfTemplate.lower() == "no" or VariantOfTemplate.lower() == "нет": #ReadyTeamplate
    alittleofhistory =  MyPermit("Напиши немного о истории "+subject,uslovia +language)
    describtionofsubject = MyPermit("Напиши определение "+subject,uslovia +language)
    describtionofpresentation = MyPermit("Напиши про представление объекта (вещи)"+subject,uslovia +language)
    objectiveoftheassignedobject = MyPermit("Напиши задачу объекта (вещи) "+subject,uslovia +language)
    conclution = MyPermit("Напиши заключение по написанным ранее четырём строкам",uslovia +language)

    '''
    alittleofhistory = input("input peace of subject's history")
    describtionofsubject= input("input describtion of subject")
    describtionofpresentation= input("input describtion of presentation")
    objectiveoftheassignedobject= input("input objective of the assigned object")
    conclution= input("input conclusion")
    '''
    #TittleSlide
    slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(slide_layout)

    left = top = Inches(0)
    height = Inches(7)
    width = Inches(10)
    #pic = slide.shapes.add_picture(img_path1, left, top, height=height, width=width)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = subject

    subtitle.text = "" + placeofstudy + "\n" \
                    "" + departmentorclass + "\n" \
                    '' + autor + '\n' \
                    "" + manager + "\n" \
                    "" + place + ", " + year + ""

    #History
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = firstTemplateslide

    content.text = (alittleofhistory)

    #Notion
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = secondTemplateslide

    content.text = (describtionofsubject)

    #RepresentationKnowledge
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = thirdTemplateslide

    content.text = (describtionofpresentation)

    #Objectives
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = fourthTemplateslide

    content.text = (objectiveoftheassignedobject)

    #Conclusion
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = fifthTemplateslide

    content.text = (conclution)



presentation.save("template.pptx")