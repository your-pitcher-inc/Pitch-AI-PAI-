from pptx import Presentation
from pptx.util import Inches
from gigachat import GigaChat
from gigachat.models import Chat, Messages, MessagesRole

def MyPermit(permit, content, giga):
    if type(permit) is str and type(content) is str:
        payload = Chat(
            messages=[
                Messages(
                    role=MessagesRole.SYSTEM,
                    content=content
                )
            ],
            temperature=0.7,
            max_tokens=100,
        )
        payload.messages.append(Messages(role=MessagesRole.USER, content=permit))
        response = giga.chat(payload)
        payload.messages.append(response.choices[0].message)
        print(response.choices[0].message.content)
        return str(response.choices[0].message.content)

def create_presentation(lang, qwsb, plos, doc, aut, mng, plc, year, vot, nos, nmos_values):
    giga = GigaChat(credentials="ZTcwNWE2ZTYtMzNmNy00YTZmLThlMDMtMGE0NjdhYWJkMTY1OmUwZmQzNjY0LTNhYmUtNGRjZC1hMDJlLTM3M2JlOTk1NjRlNg==", verify_ssl_certs=False)
    if lang == "rus":
        language = "Напиши это на русском"
        first_template_slide = "История"
        second_template_slide = "Основные понятия"
        third_template_slide = "Представление знаний"
        fourth_template_slide = "Задачи формирования баз знаний"
        fifth_template_slide = "Заключение"
    elif lang == "ing":
        language = "Write it in English"
        first_template_slide = "History"
        second_template_slide = "Basic concepts of "
        third_template_slide = "Representation knowledge"
        fourth_template_slide = "Objectives"
        fifth_template_slide = "Conclusion"
    else:
        raise ValueError("Invalid language input. Please choose 'rus' or 'ing'.")

    # Gather information for the presentation
    subject = qwsb
    place_of_study = plos
    department_or_class = doc
    author = aut
    manager = mng
    place = plc
    presentation_year = year

    slides = []

    # Create a PowerPoint presentation
    presentation = Presentation()

    # CustomTemplate
    variant_of_template = vot
    if variant_of_template.lower() == "yes" or variant_of_template.lower() == "да":
        number_of_slides = int(nos)
        for i in range(number_of_slides):
            slides.append(nmos_values[i])

        # Title Slide
        slide_layout = presentation.slide_layouts[2]
        slide = presentation.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = subject
        subtitle.text = (
            f"{place_of_study}\n{department_or_class}\n{author}\n{manager}\n{place}, {presentation_year}"
        )
        number = 0
        for i in range(number_of_slides):
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = slides[number]

            content.text = MyPermit(
                f"Write {slides[number]} notion of {subject}", language, giga
            )  # GigaChat should respond with: write [Notion of SUBJECT(SLIDE NAME)]

            number += 1

    elif variant_of_template.lower() == "no" or variant_of_template.lower() == "нет":
        alittle_of_history = MyPermit(
            f"Write a little about the history of {subject}", language, giga
        )
        description_of_subject = MyPermit(
            f"Write a definition of {subject}", language, giga
        )
        description_of_presentation = MyPermit(
            f"Write about the representation of {subject}", language, giga
        )
        objective_of_the_assigned_object = MyPermit(
            f"Write about the objectives of {subject}", language, giga
        )
        conclusion = MyPermit(
            f"Write a conclusion based on the previous four lines", language, giga
        )

        # Title Slide
        slide_layout = presentation.slide_layouts[2]
        slide = presentation.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = subject

        subtitle.text = (
            f"{place_of_study}\n{department_or_class}\n{author}\n{manager}\n{place}, {presentation_year}"
        )

        # History
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = first_template_slide

        content.text = alittle_of_history

        # Notion
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = second_template_slide

        content.text = description_of_subject

        # RepresentationKnowledge
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = third_template_slide

        content.text = description_of_presentation

        # Objectives
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = fourth_template_slide

        content.text = objective_of_the_assigned_object

        # Conclusion
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = fifth_template_slide

        content.text = conclusion

    presentation.save("template.pptx")


# Uncomment the next line to execute the presentation creation
# create_presentation("rus", "", "", "", "", "", "", "", "yes", "2", ["Slide1", "Slide2"])
