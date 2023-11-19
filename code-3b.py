from pptx import Presentation
from pptx.util import Inches
from gigachat import GigaChat
from gigachat.models import Chat, Messages, MessagesRole

def MyPermit(permit, contect):
    if type(permit) is str and type(contect) is str:
        payload = Chat(
            messages=[
                Messages(
                    role=MessagesRole.SYSTEM,
                    content=contect
                )
            ],
            temperature=0.7,
            max_tokens=100,
        )
        with GigaChat(credentials="ZTcwNWE2ZTYtMzNmNy00YTZmLThlMDMtMGE0NjdhYWJkMTY1OmUwZmQzNjY0LTNhYmUtNGRjZC1hMDJlLTM3M2JlOTk1NjRlNg==", verify_ssl_certs=False) as giga:
            payload.messages.append(Messages(role=MessagesRole.USER, content=permit))
            response = giga.chat(payload)
            payload.messages.append(response.choices[0].message)
            print(response.choices[0].message.content)
    return str(response.choices[0].message.content)

uslovia = ". Распиши это так, как будто ты делаешь доклад в очень строгой форме. Уложи это в менее 30 слов и одно предложение. Говори от третьего лица и не расписывай свои действия. Не принимай себя ни за мужчину, ни за женщину"

def create_presentation(lang, qwsb, plos, doc, aut, mng, plc, year, vot):
    # Language-specific settings
    if lang == "rus" or lang == "рус":
        language = "Напиши это на русском"
        firstTemplateslide = "История"
        secondTemplateslide = "Основные понятия"
        thirdTemplateslide = "Представление знаний"
        fourthTemplateslide = "Задачи формирования баз знаний"
        fifthTemplateslide = "Заключение"
    elif lang == "ing" or lang == "англ":
        language = "Write it in English"
        firstTemplateslide = "History"
        secondTemplateslide = "Basic concepts of "
        thirdTemplateslide = "Representation knowledge"
        fourthTemplateslide = "Objectives"
        fifthTemplateslide = "Conclusion"
    else:
        raise ValueError("Invalid language input. Please choose 'rus' or 'ing'.")

    subject = qwsb
    placeofstudy = plos
    departmentorclass = doc
    autor = aut
    manager = mng
    place = plc
    year = year

    slides = []

    presentation = Presentation()

    # CustomTemplate
    VariantOfTemplate = vot
    if VariantOfTemplate.lower() == "yes" or VariantOfTemplate.lower() == "да":
        numberofslides = int(input("How many slides do you need? Enter a number: "))
        for i in range(numberofslides):
            slides.append(input("Input name of next slide: "))

        # TitleSlide
        slide_layout = presentation.slide_layouts[2]
        slide = presentation.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = subject
        subtitle.text = (
            f"{placeofstudy}\n{departmentorclass}\n{autor}\n{manager}\n{place}, {year}"
        )
        number = 0
        for i in range(numberofslides):
            slide_layout = presentation.slide_layouts[1]
            slide = presentation.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = slides[number]

            content.text = MyPermit(
                f"Write {slides[number]} notion of {subject}", uslovia + language
            )  # GigaChat should respond with: write [Notion of SUBJECT(SLIDE NAME)]

            number += 1

    elif VariantOfTemplate.lower() == "no" or VariantOfTemplate.lower() == "нет":
        alittleofhistory = MyPermit(
            f"Write a little about the history of {subject}", uslovia + language
        )
        describtionofsubject = MyPermit(
            f"Write a definition of {subject}", uslovia + language
        )
        describtionofpresentation = MyPermit(
            f"Write about the representation of {subject}", uslovia + language
        )
        objectiveoftheassignedobject = MyPermit(
            f"Write about the objectives of {subject}", uslovia + language
        )
        conclusion = MyPermit(
            f"Write a conclusion based on the previous four lines", uslovia + language
        )

        # TitleSlide
        slide_layout = presentation.slide_layouts[2]
        slide = presentation.slides.add_slide(slide_layout)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = subject

        subtitle.text = (
            f"{placeofstudy}\n{departmentorclass}\n{autor}\n{manager}\n{place}, {year}"
        )

        # History
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = firstTemplateslide

        content.text = alittleofhistory

        # Notion
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = secondTemplateslide

        content.text = describtionofsubject

        # RepresentationKnowledge
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = thirdTemplateslide

        content.text = describtionofpresentation

        # Objectives
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = fourthTemplateslide

        content.text = objectiveoftheassignedobject

        # Conclusion
        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = fifthTemplateslide

        content.text = conclusion

    else:
        raise ValueError("Invalid template input. Please choose 'yes' or 'no'.")

    presentation.save("template.pptx")

create_presentation("rus", "Тема", "Место", "Кафедра", "Имя", "Руководитель", "Город", "2023", "да")