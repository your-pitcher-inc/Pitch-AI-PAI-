from pptx import Presentation
from pptx.util import Inches, Pt
from gigachat import GigaChat



def MyPermit(permit):
    if type(permit) is str:
        l = ""
        with GigaChat(credentials="MmEyNjJiNGUtY2UzYy00MWYxLTllNTEtOGEyZTRmMTQzNzg5OjZhZDdlMDJiLWMzYjItNGNmZi1iNDc3LTYxNWM2OGMyODg3ZQ==", verify_ssl_certs=False) as giga:
            if isEng == 1:
                l = ". Дай ответ на английском."
            print(permit + l)
            response = giga.chat(permit + l)
            print(response.choices[0].message.content)
    return str(response.choices[0].message.content)

isEng=0

img_path1 = r'C:\Users\ars\Pictures\00000.png'
img_path2 = r'C:\Users\ars\Pictures\00000.png'

#TittleList
if input("input language. ").lower() == "english":
    isEng = 1
subject= input("input subject. ")
placeofstudy= input("input your place of study. ")
departmentorclass= input("input your academic departament or academic class. ")
autor= input("input your name and surname. ")
manager= input("input your manager's name and surname. ")
place= input("input your location (city or another settlement). ")
year= input("input year of presentation. ")

slides = []

#Создаем объект Presentation
presentation = Presentation()

#CustomTemplate
VariantOfTemplate = input("Do you have your own template for this presentation? ")
if VariantOfTemplate.lower() == "yes":
    numberofslides = int(input("How many slides do you need? (Enter a number) "))
    for i in range (numberofslides):
            slides.append(input("Input name of next slide. "))

    #TittleSlide
    slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(slide_layout)

    left = top = Inches(0)
    height = Inches(7)
    width = Inches(10)
    #pic = slide.shapes.add_picture(img_path1, left, top, height=height, width=width)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = subject
    subtitle.text = "" + placeofstudy + "\n" \
                    "" + departmentorclass + "\n" \
                    '' + autor + '\n' \
                    "" + manager + "\n" \
                    "" + place + ", " + year + ""
    for i in range (numberofslides):

        number = 0

        slide_layout = presentation.slide_layouts[1]
        slide = presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = slides[number]

        content.text = (input("Input "+slides[number]+" of "+ subject)) #на это должен отвечать gigachad в виде: write [Notion of SUBJECT(НАЗВАНИЕ СЛАЙДА)]

        number+=1


elif VariantOfTemplate.lower() == "no": #ReadyTeamplate

    alittleofhistory = MyPermit("Напиши немного о истории "+subject+", как будто ты делаешь доклад. Уложи это в менее 30 слов. Не принимай себя ни за мужчину, ни за женщину")
    describtionofsubject = MyPermit("Напиши определение "+subject+", как будто ты делаешь доклад. Уложи это в менее 30 слов. Не принимай себя ни за мужчину, ни за женщину")
    describtionofpresentation = MyPermit("Напиши про представление знаний "+subject+", как будто ты делаешь доклад. Уложи это в менее 30 слов. Не принимай себя ни за мужчину, ни за женщину")
    objectiveoftheassignedobject = MyPermit("Напиши задачу формирования баз знаний "+subject+", как будто ты делаешь доклад. Уложи это в менее 30 слов. Не принимай себя ни за мужчину, ни за женщину")
    conclution = MyPermit("Напиши заключение по написанным ранее четырём строкам. Уложи это в менее 30 слов. Не принимай себя ни за мужчину, ни за женщину")

    '''
    alittleofhistory = input("input peace of subject's history")
    describtionofsubject= input("input describtion of subject")
    describtionofpresentation= input("input describtion of presentation")
    objectiveoftheassignedobject= input("input objective of the assigned object")
    conclution= input("input conclusion")
    '''
    #TittleSlide
    slide_layout = presentation.slide_layouts[2]
    slide = presentation.slides.add_slide(slide_layout)

    left = top = Inches(0)
    height = Inches(7)
    width = Inches(10)
    #pic = slide.shapes.add_picture(img_path1, left, top, height=height, width=width)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = subject

    subtitle.text = "" + placeofstudy + "\n" \
                    "" + departmentorclass + "\n" \
                    '' + autor + '\n' \
                    "" + manager + "\n" \
                    "" + place + ", " + year + ""

    #History
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = 'Notion of ' + subject

    content.text = (alittleofhistory)

    #Notion
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = 'Basic concepts of ' + subject

    content.text = (describtionofsubject)

    #RepresentationKnowledge
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = 'Representation knowledge'

    content.text = (describtionofpresentation)

    #Objectives
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = 'Objectives'

    content.text = (objectiveoftheassignedobject)

    #Conclusion
    slide_layout = presentation.slide_layouts[1]
    slide = presentation.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = 'Conclusion'

    content.text = (conclution)

presentation.save("template.pptx")




'''
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

left = top =  Inches(0)
height = Inches(7)
width = Inches(10)
pic = slide.shapes.add_picture(img_path2, left, top, height=height, width = width)
'''


